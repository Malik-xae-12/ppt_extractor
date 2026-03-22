import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import base64
from typing import List
from app.modules.extract.schema import ExtractedSlide, ExtractionResponse
from markitdown import MarkItDown

def encode_image(image_bytes: bytes) -> str:
    """Encode image bytes to base64."""
    return f"data:image/png;base64,{base64.b64encode(image_bytes).decode('utf-8')}"

def extract_content(file_path: str, filename: str) -> ExtractionResponse:
    # 1. Use markitdown for overall translation to markdown
    md = MarkItDown()
    try:
        result = md.convert(file_path)
        full_markdown = result.text_content
    except Exception as e:
        full_markdown = f"Error generating markdown: {e}"

    # 2. Use python-pptx for detailed slide-by-slide extraction
    prs = Presentation(file_path)
    slides_data: List[ExtractedSlide] = []

    for idx, slide in enumerate(prs.slides):
        text_content = ""
        images = []
        tables = []
        architecture_diagrams = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_content += shape.text + "\n\n"
            
            # Extract standard pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob
                    images.append(encode_image(img_bytes))
                    # Pictures might also be architecture diagrams if they are large or complex
                    # but we keep it simple here.
                except Exception:
                    pass
                
            # Treat grouped shapes or SmartArt as potential architecture diagrams
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                diagram_text = ""
                for s in shape.shapes:
                    if s.has_text_frame:
                        diagram_text += s.text + " "
                architecture_diagrams.append(f"Diagram (Group Shape) containing text: {diagram_text.strip()}")
            
            # Extract tables
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                tables.append(table_data)
        
        slides_data.append(
            ExtractedSlide(
                slide_number=idx + 1,
                text_content=text_content.strip(),
                markdown_content="",  # we can attempt to segment markdown if needed
                images=images,
                tables=tables,
                architecture_diagrams=architecture_diagrams
            )
        )
    
    return ExtractionResponse(
        filename=filename,
        full_markdown=full_markdown,
        slides=slides_data
    )
